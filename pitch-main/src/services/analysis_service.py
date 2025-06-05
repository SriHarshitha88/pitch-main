import os
import json
from sqlalchemy.orm import Session
from datetime import datetime # Import datetime
from .. import database, models, config # Import config
# Assuming you have an AI client initialized elsewhere or will initialize it here
# import openai # Import openai
from groq import Groq # Import Groq
# Import libraries for reading different file types
from pptx import Presentation # For .pptx
from pypdf import PdfReader # For .pdf

print("analysis_service.py is being loaded") # Added print statement

def read_file_content(file_path: str) -> str:
    """Read content from PDF or PPTX file."""
    print(f"Reading file: {file_path}")
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pdf':
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    elif file_ext in ['.pptx', '.ppt']:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

def perform_analysis(analysis_id: int, file_path: str):
    """Perform AI analysis on a pitch deck."""
    print(f"Starting analysis for job {analysis_id}")
    print(f"Using Groq API key: {config.GROQ_API_KEY[:5]}...{config.GROQ_API_KEY[-4:]}")
    
    # Get database session
    session = Session(database.engine)
    analysis = session.query(models.Analysis).get(analysis_id)
    if analysis is None:
        raise ValueError(f"Analysis with id {analysis_id} not found")

    try:
        # Update status to processing
        analysis.status = "processing"
        session.commit()

        # Read file content
        file_content = read_file_content(file_path)
        print(f"Successfully read file content, length: {len(file_content)}")

        # Prepare the prompt
        prompt = """Analyze the following pitch deck content and provide a detailed analysis in JSON format.
The response should include:
1. overall_score (0-100)
2. pitch_analysis (clarity, storytelling, value proposition)
3. market_research (market size, competition, growth potential)
4. financial_analysis (revenue model, projections, funding needs)
5. generated_report (detailed analysis in markdown format)

Pitch deck content:
{content}

Respond with a valid JSON object containing these sections."""

        # Initialize Groq client
        print("Initializing Groq client...")
        try:
            client = Groq(api_key=config.GROQ_API_KEY)
            print("Successfully initialized Groq client")
        except Exception as client_error:
            print(f"Error initializing Groq client: {str(client_error)}")
            raise

        # Make API call
        print("Preparing to make Groq API call...")
        try:
            print("Creating chat completion request...")
            response = client.chat.completions.create(
                model="llama3-8b-8192",  # Updated to use supported model
                response_format={"type": "json_object"},
                messages=[
                    {"role": "system", "content": "You are a pitch deck analysis AI that provides detailed, structured analysis in JSON format."},
                    {"role": "user", "content": prompt.format(content=file_content)}
                ],
                temperature=0.7,
                max_tokens=4096
            )
            print("Successfully made Groq API call")
        except Exception as api_error:
            print(f"Error during Groq API call: {str(api_error)}")
            raise
        
        print("Received response from Groq API")
        
        if not response:
            print("Response object is None")
            raise ValueError("Groq API returned None response")
        
        if not hasattr(response, 'choices'):
            print(f"Response object has no 'choices' attribute. Response type: {type(response)}")
            raise ValueError("Groq API response has no 'choices' attribute")
            
        if not response.choices:
            print("Response.choices is empty")
            raise ValueError("Groq API returned empty choices")

        # Extract and parse the response
        print("Extracting response content...")
        ai_response_text = response.choices[0].message.content
        print(f"Raw AI response: {ai_response_text[:200]}...")  # Print first 200 chars
        
        try:
            print("Parsing JSON response...")
            analysis_result = json.loads(ai_response_text)
            print("Successfully parsed JSON response")
        except json.JSONDecodeError as e:
            print(f"Error parsing JSON response: {e}")
            raise ValueError(f"Invalid JSON response from AI: {e}")

        # Create AnalysisResult record
        print("Creating AnalysisResult record...")
        result = models.AnalysisResult(
            analysis_id=analysis.id,
            overall_score=analysis_result.get('overall_score', 0),
            pitch_analysis=analysis_result.get('pitch_analysis', {}),
            market_research=analysis_result.get('market_research', {}),
            financial_analysis=analysis_result.get('financial_analysis', {}),
            website_analysis=analysis_result.get('website_analysis', {}),
            investment_strategy=analysis_result.get('investment_strategy', {}),
            due_diligence=analysis_result.get('due_diligence', {})
        )
        session.add(result)

        # Update analysis record
        print("Updating analysis record...")
        analysis.status = "completed"
        analysis.result = analysis_result  # Store raw result
        analysis.completed_at = datetime.utcnow()
        session.commit()
        
        print(f"Analysis completed successfully for job {analysis_id}")
        return analysis_result

    except Exception as e:
        print(f"Error during analysis for job {analysis_id}: {str(e)}")
        analysis.status = "failed"
        analysis.error = str(e)
        session.commit()
        raise

    finally:
        session.close()

    # ... existing code ... 